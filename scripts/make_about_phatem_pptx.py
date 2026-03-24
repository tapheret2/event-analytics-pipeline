from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

from PIL import Image, ImageDraw, ImageFont


OUT_DIR = Path(r"C:\Users\ADMIN\.openclaw\workspace\out")
ASSETS_DIR = OUT_DIR / "assets"
OUT_DIR.mkdir(parents=True, exist_ok=True)
ASSETS_DIR.mkdir(parents=True, exist_ok=True)

PPTX_PATH = OUT_DIR / "PhatEm_GioiThieu.pptx"


def make_mascot_png(path: Path):
    w, h = 900, 900
    img = Image.new("RGB", (w, h), (16, 24, 40))
    d = ImageDraw.Draw(img)

    # gradient-ish stripes
    for i in range(0, h, 6):
        c = (16 + i // 30, 24 + i // 40, 40 + i // 18)
        d.rectangle([0, i, w, i + 6], fill=c)

    # abstract "assistant" badge
    d.ellipse([120, 120, 780, 780], outline=(240, 248, 255), width=10)
    d.ellipse([165, 165, 735, 735], outline=(64, 196, 255), width=6)

    # simple "lobster" claws as shapes (abstract)
    d.rounded_rectangle([250, 420, 410, 560], radius=40, fill=(255, 99, 71))
    d.rounded_rectangle([490, 420, 650, 560], radius=40, fill=(255, 99, 71))
    d.ellipse([285, 445, 375, 535], fill=(16, 24, 40))
    d.ellipse([525, 445, 615, 535], fill=(16, 24, 40))

    # center dot
    d.ellipse([430, 430, 470, 470], fill=(255, 215, 0))

    # text
    try:
        font_big = ImageFont.truetype("arial.ttf", 72)
        font_mid = ImageFont.truetype("arial.ttf", 36)
    except Exception:
        font_big = ImageFont.load_default()
        font_mid = ImageFont.load_default()

    d.text((w // 2, 240), "PHÁT EM", font=font_big, fill=(240, 248, 255), anchor="mm")
    d.text((w // 2, 305), "AI Assistant (OpenClaw)", font=font_mid, fill=(180, 210, 255), anchor="mm")
    d.text((w // 2, 640), "Gọn – Rõ – Làm được", font=font_mid, fill=(220, 240, 255), anchor="mm")

    img.save(path)


def set_slide_bg(slide, rgb: tuple[int, int, int]):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)


def add_title(slide, title: str, subtitle: str | None = None):
    left, top, width, height = Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(240, 248, 255)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(22)
        p2.font.color.rgb = RGBColor(180, 210, 255)
        p2.space_before = Pt(6)


def add_bullets(slide, header: str, bullets: list[str], left=0.9, top=2.1, width=7.2, height=4.8):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()

    p0 = tf.paragraphs[0]
    p0.text = header
    p0.font.size = Pt(28)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(240, 248, 255)

    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(220, 235, 255)
        p.space_before = Pt(6)


def add_footer(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.7), Inches(0.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(160, 190, 220)
    p.alignment = PP_ALIGN.RIGHT


def main():
    mascot = ASSETS_DIR / "phatem.png"
    make_mascot_png(mascot)

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # theme colors
    BG = (10, 18, 32)
    ACCENT = (64, 196, 255)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_title(slide, "Phát Em", "Trợ lý AI của Phát — làm việc theo kiểu: rõ việc → làm được → báo lại")
    slide.shapes.add_picture(str(mascot), Inches(9.0), Inches(1.6), height=Inches(4.8))
    # accent line
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.8), Inches(1.85), Inches(7.8), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*ACCENT)
    line.line.fill.background()
    add_footer(slide, "OpenClaw • deck auto-generated")

    # Slide 2: I can help with
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_title(slide, "Mình làm được gì?")
    add_bullets(
        slide,
        "Những việc mình hay giúp Phát:",
        [
            "Theo dõi deadline (Moodle) → sync Notion → gửi digest 6:00 & 18:00",
            "Lộ trình học (DE): SQL, data modeling, pipeline, best practices",
            "Tự động hóa lặt vặt: email, file, script, checklist",
            "Giải thích vấn đề: rõ ràng, có ví dụ, có bước làm",
        ],
    )
    # small icon blocks
    for i, txt in enumerate(["Browser", "Notion API", "Gmail", "SQL"]):
        x = 9.0
        y = 2.2 + i * 1.2
        r = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.6), Inches(0.85))
        r.fill.solid(); r.fill.fore_color.rgb = RGBColor(20, 34, 60)
        r.line.color.rgb = RGBColor(*ACCENT)
        t = r.text_frame
        t.text = txt
        t.paragraphs[0].font.size = Pt(18)
        t.paragraphs[0].font.color.rgb = RGBColor(230, 245, 255)
    add_footer(slide, "Phát Em")

    # Slide 3: Principles
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_title(slide, "Cách mình làm việc")
    add_bullets(
        slide,
        "Nguyên tắc:",
        [
            "Không vòng vo: nói thẳng cái gì làm được/cái gì thiếu",
            "Tự động hóa nhưng có kiểm soát: cái gì nhạy cảm sẽ hỏi trước",
            "Ưu tiên bền: ít click UI, dùng overview/index, parse dữ liệu ổn định",
            "Ghi nhớ đúng: nếu sai thì sửa memory ngay",
        ],
    )
    slide.shapes.add_picture(str(mascot), Inches(9.3), Inches(2.3), height=Inches(3.9))
    add_footer(slide, "Phát Em")

    # Slide 4: Example (Moodle -> Notion)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_title(slide, "Ví dụ thật: Moodle → Notion")
    add_bullets(
        slide,
        "Case: 2 môn trên courses.hcmus.edu.vn",
        [
            "Check assignment/quiz overview",
            "Upsert Notion task (English), Category=Class, Priority=Medium",
            "Gửi digest theo Notion thay vì spam từng bài",
        ],
    )
    # highlight card
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.55), Inches(11.8), Inches(1.35))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(20, 34, 60)
    card.line.color.rgb = RGBColor(*ACCENT)
    tf = card.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "MTH10131 - BTVN#3  |  Due: 2026-03-24 12:00  |  Status: Not started"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(230, 245, 255)
    p2 = tf.add_paragraph()
    p2.text = "Description: exercises 2.6, 2.7, 2.10, 2.11, 2.14, 2.15, 2.16 + Moodle link"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(180, 210, 255)
    add_footer(slide, "Phát Em")

    # Slide 5: What I need from you
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_title(slide, "Để mình chạy mượt, mình cần gì từ Phát?")
    add_bullets(
        slide,
        "Chỉ 3 thứ:",
        [
            "Link/nguồn dữ liệu rõ ràng (Moodle, Drive, GitHub…)",
            "Quy tắc giao tiếp/nhắc nhở (giờ báo, format, mức chi tiết)",
            "Token/API khi cần (Notion token, v.v.) — và rotate khi lỡ lộ",
        ],
    )
    # decorative shapes
    for i in range(3):
        c = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.3 + i*1.1), Inches(2.6), Inches(0.8), Inches(0.8))
        c.fill.solid(); c.fill.fore_color.rgb = RGBColor(64, 196, 255) if i==0 else RGBColor(255, 99, 71) if i==1 else RGBColor(255, 215, 0)
        c.line.fill.background()
    add_footer(slide, "Phát Em")

    # Slide 6: Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_title(slide, "Ok, giờ mình làm gì tiếp?", "Gợi ý: DE roadmap • project • sync Notion • đọc deadline")
    slide.shapes.add_picture(str(mascot), Inches(9.0), Inches(2.0), height=Inches(4.6))
    # CTA box
    cta = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(4.0), Inches(7.8), Inches(1.6))
    cta.fill.solid(); cta.fill.fore_color.rgb = RGBColor(20, 34, 60)
    cta.line.color.rgb = RGBColor(*ACCENT)
    tf = cta.text_frame
    tf.text = "Nhắn mình kiểu: \n• 'làm roadmap DE 8 tuần'\n• 'tạo project ETL + Docker'\n• 'tối nay check Notion digest'"
    for idx, p in enumerate(tf.paragraphs):
        p.font.color.rgb = RGBColor(230, 245, 255)
        p.font.size = Pt(18 if idx==0 else 18)
    add_footer(slide, "Phát Em")

    prs.save(PPTX_PATH)
    print(str(PPTX_PATH))


if __name__ == "__main__":
    main()
